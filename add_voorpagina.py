#!/usr/bin/env python3
"""Voeg 1 voorpagina toe aan bestaande pivot-presentatie. Raakt andere slides niet aan."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from copy import deepcopy

DARK = RGBColor(0x2D, 0x29, 0x38)
BG_SLIDE = RGBColor(0xF8, 0xF5, 0xEF)
REFINE = RGBColor(0xF5, 0x9E, 0x0B)
MUTED = RGBColor(0x6B, 0x61, 0x73)

PATH = "/Users/jr/Library/CloudStorage/OneDrive-Zonneplan/School/Innovatiemanagement/public/presentatie-pivot.pptx"
prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb


# ============ VOORPAGINA ============
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Achtergrond
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid(); bg.fill.fore_color.rgb = BG_SLIDE
slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)

# Decoratieve accent-balk links
accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.2), Inches(0.45), Inches(3.1))
accent_bar.line.fill.background()
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = REFINE

# Tag boven
tag_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(0.95), Inches(0.5), Inches(0.08))
tag_bar.line.fill.background()
tag_bar.fill.solid(); tag_bar.fill.fore_color.rgb = REFINE
add_text(slide, Inches(1.65), Inches(0.82), Inches(8), Inches(0.4),
         "INNOVATIEMANAGEMENT · EDIM.21", size=14, bold=True, color=REFINE)

# Hoofdtitel
add_text(slide, Inches(1.0), Inches(2.2), Inches(11.5), Inches(1.2),
         "Aurora", size=72, bold=True, color=DARK)
add_text(slide, Inches(1.0), Inches(3.15), Inches(11.5), Inches(1.0),
         "Beyond Paper", size=52, bold=True, color=REFINE)

# Subtitel / onderwerp
add_text(slide, Inches(1.0), Inches(4.35), Inches(11.5), Inches(0.8),
         "De pivot: van natte handtekening naar gekoppelde identiteit",
         size=22, color=DARK)

# Lijn
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.45), Inches(3.5), Inches(0.03))
divider.line.fill.background()
divider.fill.solid(); divider.fill.fore_color.rgb = DARK

# Team
add_text(slide, Inches(1.0), Inches(5.65), Inches(11), Inches(0.4),
         "Anke Maassen van den Brink  ·  Lianne van Os  ·  Jr Bouwman",
         size=16, color=DARK)
add_text(slide, Inches(1.0), Inches(6.1), Inches(11), Inches(0.4),
         "Technische Bedrijfskunde  ·  Windesheim  ·  2025–2026",
         size=13, color=MUTED)

# Casus-label rechtsonder
add_text(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
         "Casus: Aurora Writing Instruments Group",
         size=11, color=MUTED, align=PP_ALIGN.RIGHT)


# ============ Verplaats nieuw toegevoegde slide naar positie 0 ============
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
# laatst toegevoegde = ons voorblad
new_slide = slides_list[-1]
xml_slides.remove(new_slide)
xml_slides.insert(0, new_slide)

prs.save(PATH)
print(f"Voorpagina toegevoegd als slide 1. Saved: {PATH}")
