#!/usr/bin/env python3
"""Pivot presentatie: hoogover + strategische/financiële toelichting. Max 3 slides, weinig tekst."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Huisstijl
DARK = RGBColor(0x2D, 0x29, 0x38)
BG = RGBColor(0x0F, 0x0E, 0x1A)
BG_SLIDE = RGBColor(0xF8, 0xF5, 0xEF)  # cream
REFINE = RGBColor(0xF5, 0x9E, 0x0B)
REFINE_LIGHT = RGBColor(0xFE, 0xF3, 0xE7)
TEXT = RGBColor(0x2D, 0x29, 0x38)
MUTED = RGBColor(0x6B, 0x61, 0x73)
CORAL = RGBColor(0xF4, 0x72, 0x64)
BG_RISK = RGBColor(0xFD, 0xEC, 0xEA)

HOOGOVER_PNG = "/tmp/hoogover-export.png"
VPC_PNG = "/tmp/vpc-export.png"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


def add_bg(slide, color=BG_SLIDE):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element); slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return tb


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Inches(0.08), color=REFINE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = color


def add_card(slide, left, top, width, height, title, bullets, *, accent=REFINE, bg_color=None, text_color=TEXT):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.line.color.rgb = accent
    card.line.width = Pt(1.75)
    card.fill.solid(); card.fill.fore_color.rgb = bg_color if bg_color else RGBColor(0xFF, 0xFF, 0xFF)

    # Accent bar
    add_accent_bar(slide, left + Inches(0.3), top + Inches(0.3), width=Inches(0.45), height=Inches(0.06), color=accent)

    # Title
    add_text(slide, left + Inches(0.3), top + Inches(0.42), width - Inches(0.5), Inches(0.5),
             title, size=18, bold=True, color=text_color, font="Calibri")

    # Bullets
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.0), width - Inches(0.5), height - Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run(); run.text = "•  " + b
        run.font.size = Pt(14); run.font.color.rgb = MUTED; run.font.name = "Calibri"


# ============ SLIDE 1: Hoogover flowchart ============
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1, BG_SLIDE)

# Tag
add_accent_bar(s1, Inches(0.6), Inches(0.45), width=Inches(0.4), height=Inches(0.06))
add_text(s1, Inches(1.1), Inches(0.36), Inches(6), Inches(0.3),
         "PIVOT · AURORA", size=12, bold=True, color=REFINE)

# Title
add_text(s1, Inches(0.6), Inches(0.62), Inches(11), Inches(1.0),
         "Van natte handtekening naar gekoppelde identiteit",
         size=34, bold=True, color=TEXT)

add_text(s1, Inches(0.6), Inches(1.55), Inches(11), Inches(0.4),
         "Het papier en de pen blijven. Het bewijs verschuift.",
         size=17, color=MUTED)

# Hoogover image
img_left = Inches(0.6)
img_top = Inches(2.1)
img_w = Inches(12.1)
s1.shapes.add_picture(HOOGOVER_PNG, img_left, img_top, width=img_w)

# Footer
add_text(s1, Inches(0.6), Inches(7.15), Inches(8), Inches(0.25),
         "Anke · Lianne · Jr  ·  EDIM.21", size=11, color=MUTED)


# ============ SLIDE 2: Strategisch + technisch ============
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2, BG_SLIDE)

add_accent_bar(s2, Inches(0.6), Inches(0.45), width=Inches(0.4), height=Inches(0.06))
add_text(s2, Inches(1.1), Inches(0.36), Inches(6), Inches(0.3),
         "ANALYSE · DEEL 1", size=12, bold=True, color=REFINE)

add_text(s2, Inches(0.6), Inches(0.62), Inches(11), Inches(0.7),
         "Strategisch & technisch", size=32, bold=True, color=TEXT)

# 2x2 cards
card_w = Inches(6.0)
card_h = Inches(2.85)
col1, col2 = Inches(0.6), Inches(6.75)
row1, row2 = Inches(1.65), Inches(4.55)

add_card(s2, col1, row1, card_w, card_h,
         "Strategische fit",
         ["125+ jaar autoriteit rond ondertekenen",
          "Vertrouwen en gezag blijven",
          "Van tool-maker naar procespartner"])

add_card(s2, col2, row1, card_w, card_h,
         "Kerncompetenties",
         ["Domeinkennis over ondertekenen",
          "Merkautoriteit bij notaris en bank",
          "Bestaand distributienetwerk"])

add_card(s2, col1, row2, card_w, card_h,
         "Technische fit",
         ["Aurora: kennismodel + merk",
          "Techpartner: infrastructuur",
          "Hergebruik ID-scan, NFC, PIN"])

add_card(s2, col2, row2, card_w, card_h,
         "Concurrentievoordeel",
         ["DocuSign/PKI: volledig digitaal",
          "DigiD: overheidsgebonden",
          "Witte plek: hybride papier + digitaal bewijs"])


# ============ SLIDE 3: Customer value + financieel ============
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3, BG_SLIDE)

add_accent_bar(s3, Inches(0.6), Inches(0.45), width=Inches(0.4), height=Inches(0.06))
add_text(s3, Inches(1.1), Inches(0.36), Inches(6), Inches(0.3),
         "ANALYSE · DEEL 2", size=12, bold=True, color=REFINE)

add_text(s3, Inches(0.6), Inches(0.62), Inches(11), Inches(0.7),
         "Waarde & financieel", size=32, bold=True, color=TEXT)

# 3 columns
card_w3 = Inches(4.05)
card_h3 = Inches(5.4)
gap = Inches(0.2)
col_a = Inches(0.5)
col_b = col_a + card_w3 + gap
col_c = col_b + card_w3 + gap
row_y = Inches(1.65)

add_card(s3, col_a, row_y, card_w3, card_h3,
         "Customer value",
         ["Notaris: rechtsgeldig zonder extra stap",
          "Bank: fraudepreventie bij contracten",
          "Gemeente: schaalbare ondertekening",
          "Gebruiker: vertrouwde handtekening, bewijsbare identiteit"])

add_card(s3, col_b, row_y, card_w3, card_h3,
         "Financiële kansen",
         ["Licentiemodel op kennismodel",
          "Hardware-omzet verificatiestations",
          "Training en consultancy",
          "NL-markt: duizenden notarissen, banken, gemeenten"])

add_card(s3, col_c, row_y, card_w3, card_h3,
         "Financiële risico's",
         ["Ontwikkelkosten techpartner buiten grip",
          "Adoptietijd lang (compliance)",
          "Afhankelijk van één techpartner",
          "Digitaal-only concurrentie drukt marges"],
         accent=CORAL, bg_color=BG_RISK)


# ============ SLIDE 4: Value Proposition Canvas (native shapes) ============
from pptx.oxml.ns import qn
from lxml import etree

s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4, BG_SLIDE)

add_accent_bar(s4, Inches(0.6), Inches(0.45), width=Inches(0.4), height=Inches(0.06))
add_text(s4, Inches(1.1), Inches(0.36), Inches(8), Inches(0.3),
         "VALUE PROPOSITION CANVAS", size=12, bold=True, color=REFINE)
add_text(s4, Inches(0.6), Inches(0.62), Inches(11), Inches(0.7),
         "Waar Aurora's aanbod de klant raakt", size=30, bold=True, color=TEXT)


def add_line(slide, x1, y1, x2, y2, color=DARK, weight=1.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_section_title(slide, left, top, width, text, *, align=PP_ALIGN.LEFT, color=TEXT, size=13):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    r.font.name = "Calibri"
    rPr = r._r.get_or_add_rPr()
    rPr.set("spc", "150")
    return tb


def add_bullets(slide, left, top, width, height, bullets, *, size=12, color=TEXT, dot=REFINE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(3)
        # coloured bullet
        rb = p.add_run(); rb.text = "•  "
        rb.font.size = Pt(size); rb.font.color.rgb = dot; rb.font.bold = True; rb.font.name = "Calibri"
        # text
        rt = p.add_run(); rt.text = b
        rt.font.size = Pt(size); rt.font.color.rgb = color; rt.font.name = "Calibri"
    return tb


# ============ VALUE MAP square ============
sq_left = Inches(0.45); sq_top = Inches(1.75); sq_w = Inches(6.1); sq_h = Inches(5.4)

# label boven square
add_text(s4, sq_left, sq_top - Inches(0.4), Inches(3), Inches(0.3),
         "VALUE MAP", size=14, bold=True, color=REFINE)
add_text(s4, sq_left + sq_w - Inches(3), sq_top - Inches(0.38), Inches(3), Inches(0.28),
         "Aurora + techpartner", size=11, color=MUTED, align=PP_ALIGN.RIGHT)

# outer square
square = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, sq_left, sq_top, sq_w, sq_h)
square.fill.solid(); square.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
square.line.color.rgb = DARK; square.line.width = Pt(2)

# diagonal lines from right-middle
rm_x = sq_left + sq_w
rm_y = sq_top + sq_h / 2
add_line(s4, rm_x, rm_y, sq_left, sq_top, DARK, 1.5)
add_line(s4, rm_x, rm_y, sq_left, sq_top + sq_h, DARK, 1.5)

# Section titles (inside each triangle)
add_section_title(s4, sq_left + Inches(0.2), sq_top + sq_h / 2 - Inches(0.85),
                  Inches(3.5), "PRODUCTS & SERVICES", size=14)
add_section_title(s4, sq_left + sq_w - Inches(3.0), sq_top + Inches(0.25),
                  Inches(2.8), "GAIN CREATORS", align=PP_ALIGN.RIGHT, size=14)
add_section_title(s4, sq_left + sq_w - Inches(3.0), sq_top + sq_h - Inches(0.55),
                  Inches(2.8), "PAIN RELIEVERS", align=PP_ALIGN.RIGHT, size=14)

# PRODUCTS & SERVICES content (left triangle, center area)
add_bullets(s4, sq_left + Inches(0.2), sq_top + sq_h / 2 - Inches(0.35),
            Inches(4.0), Inches(1.8),
            ["Kennismodel voor authenticiteit",
             "Verificatiestation: ID + PIN",
             "Koppeling handtekening ↔ code",
             "Implementatie-ondersteuning"],
            size=13)

# GAIN CREATORS content (top-right triangle)
add_bullets(s4, sq_left + sq_w - Inches(3.1), sq_top + Inches(0.65),
            Inches(2.9), Inches(1.4),
            ["Juridisch waterdichte audittrail",
             "Schaalbaar over transacties",
             "Compliance-klaar",
             "Domeinkennis geeft gezag"],
            size=11, align=PP_ALIGN.RIGHT)

# PAIN RELIEVERS content (bottom-right triangle)
add_bullets(s4, sq_left + sq_w - Inches(3.1), sq_top + sq_h - Inches(1.95),
            Inches(2.9), Inches(1.4),
            ["Objectieve ID-check",
             "Gescande handtekening rechtsgeldig",
             "Fraudesignalen direct zichtbaar",
             "Digitaal doorzoekbaar archief"],
            size=11, align=PP_ALIGN.RIGHT)


# ============ CUSTOMER PROFILE circle ============
cir_left = Inches(7.0); cir_top = Inches(1.75); cir_d = Inches(5.4)

# label boven circle
add_text(s4, cir_left, cir_top - Inches(0.4), Inches(3.5), Inches(0.3),
         "CUSTOMER PROFILE", size=14, bold=True, color=DARK)
add_text(s4, cir_left + cir_d - Inches(3), cir_top - Inches(0.38), Inches(3), Inches(0.28),
         "Notaris, bank, gemeente", size=11, color=MUTED, align=PP_ALIGN.RIGHT)

# outer circle
circle = s4.shapes.add_shape(MSO_SHAPE.OVAL, cir_left, cir_top, cir_d, cir_d)
circle.fill.solid(); circle.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
circle.line.color.rgb = DARK; circle.line.width = Pt(2)

# 3 radii: top, left, bottom
cx = cir_left + cir_d / 2
cy = cir_top + cir_d / 2
add_line(s4, cx, cy, cx, cir_top, DARK, 1.5)
add_line(s4, cx, cy, cir_left, cy, DARK, 1.5)
add_line(s4, cx, cy, cx, cir_top + cir_d, DARK, 1.5)

# Section titles (inside each pie slice)
add_section_title(s4, cir_left + Inches(0.6), cir_top + Inches(1.3),
                  Inches(1.8), "GAINS", color=DARK, size=14)
add_section_title(s4, cir_left + Inches(0.6), cir_top + cir_d - Inches(1.8),
                  Inches(1.8), "PAINS", color=DARK, size=14)
add_section_title(s4, cir_left + cir_d - Inches(2.4), cir_top + Inches(1.3),
                  Inches(2.2), "JOB-TO-BE-DONE", align=PP_ALIGN.RIGHT, color=DARK, size=14)

# GAINS content (top-left quarter)
add_bullets(s4, cir_left + Inches(0.6), cir_top + Inches(1.75),
            Inches(2.2), Inches(1.1),
            ["Juridische zekerheid",
             "Sneller proces",
             "Minder administratie",
             "Vertrouwen eindklant"],
            size=11, dot=DARK)

# PAINS content (bottom-left quarter)
add_bullets(s4, cir_left + Inches(0.6), cir_top + cir_d - Inches(1.4),
            Inches(2.2), Inches(1.1),
            ["Vervalsing lastig bewijsbaar",
             "Visuele ID-check foutgevoelig",
             "Scan niet rechtsgeldig",
             "Fraude komt laat aan licht"],
            size=11, dot=DARK)

# JOB-TO-BE-DONE content (right half)
add_bullets(s4, cir_left + cir_d / 2 + Inches(0.2), cir_top + cir_d / 2 - Inches(0.9),
            Inches(2.7), Inches(1.9),
            ["Rechtsgeldig laten ondertekenen",
             "Juiste persoon tekent",
             "Audittrail voor compliance",
             "Klanten snel bedienen"],
            size=12, dot=DARK)

# Save
output = "/Users/jr/Library/CloudStorage/OneDrive-Zonneplan/School/Innovatiemanagement/public/presentatie-pivot.pptx"
prs.save(output)
print(f"Saved: {output}")
