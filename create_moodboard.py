#!/usr/bin/env python3
"""
Concept Collection / Visual Moodboard: De Analoge Security Expert
V3: WARRIG, door elkaar, wild divergent
Niet alleen een pen — alle mogelijke vormen van authenticatie/verificatie
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os, random

random.seed(42)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
slide = prs.slides.add_slide(prs.slide_layouts[6])

IMG = "/Users/jr/Library/CloudStorage/OneDrive-Zonneplan/School/Innovatiemanagement/public/images/moodboard"
RESEARCH = "/Users/jr/Library/CloudStorage/OneDrive-Zonneplan/School/Innovatiemanagement/public/images/research"

# Colors
GOLD = RGBColor(0xD4, 0xA5, 0x37)
CREAM = RGBColor(0xF5, 0xF0, 0xE1)
DARK = RGBColor(0x06, 0x06, 0x0C)
CYAN = RGBColor(0x00, 0xE5, 0xFF)
PINK = RGBColor(0xFF, 0x10, 0x80)
LIME = RGBColor(0xA0, 0xFF, 0x20)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
PURPLE = RGBColor(0xBB, 0x55, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xFF, 0x22, 0x22)
YELLOW = RGBColor(0xFF, 0xEB, 0x3B)

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(16), Inches(9))
bg.fill.solid(); bg.fill.fore_color.rgb = DARK; bg.line.fill.background()

def img(name, l, t, w, h, rot=0):
    path = os.path.join(IMG, name) if "/" not in name else name
    if not os.path.exists(path): return None
    pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    if rot: pic.rotation = rot
    return pic

def tag(l, t, text, color, sz=8, w=1.8, rot=0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.28))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    if rot: s.rotation = rot
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = DARK; p.font.bold = True; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return s

def lbl(l, t, text, sz=8, color=WHITE, w=2.0, rot=0):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.6))
    tf = tx.text_frame; tf.word_wrap = True
    if rot: tx.rotation = rot
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz)
    p.font.color.rgb = color; p.font.bold = True; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return tx

def circle(l, t, sz, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(sz), Inches(sz))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

# ═══════════════════════════════════════════════════
# IMAGES — WARRIG, overlappend, gekanteld, wild
# ═══════════════════════════════════════════════════

# --- Cluster links-boven: LICHAAM ALS SLEUTEL ---
img("eye-scan.jpg",         0.1, 0.1, 3.4, 2.4, -3)
img("fingerprint.jpg",      2.6, 0.5, 2.5, 1.9, 5)
img("baby-feet.jpg",        0.3, 2.2, 2.2, 1.7, -6)
img("dna-helix.jpg",        4.6, 0.0, 2.4, 2.0, 2)

# --- Cluster midden-boven: DINGEN DIE WE AL HEBBEN ---
img("smartwatch.jpg",       6.8, 0.2, 2.4, 1.9, -4)
img("ring.jpg",             8.9, 0.0, 2.3, 2.1, 7)
img("swiss-watch.jpg",      9.0, 1.8, 2.0, 1.6, -8)

# --- Cluster rechts-boven: PATRONEN & RITME ---
img("dancing.jpg",          11.3, 0.1, 2.5, 1.9, 4)
img("chess.jpg",            13.2, 0.3, 2.7, 2.0, -5)
img("constellation.jpg",    11.0, 1.8, 2.2, 1.7, 6)

# --- Cluster links-midden: OUDE VORMEN VAN BEWIJS ---
img("wax-seal.jpg",         0.0, 3.8, 2.6, 2.0, 5)
img("door-knocker.jpg",     2.3, 3.5, 2.0, 1.6, -7)
img("lock.jpg",             0.5, 5.5, 2.3, 1.8, 3)

# --- CENTRAAL GROOT: security expert + wild eromheen ---
img(f"{RESEARCH}/denkrichting1-security-expert.jpeg", 4.2, 2.5, 4.5, 3.2, 0)
img("vault.jpg",            3.5, 5.2, 2.4, 1.8, -4)
img("circuit.jpg",          8.2, 4.8, 2.5, 1.9, 5)

# --- Cluster rechts-midden: VERRASSENDE VORMEN ---
img("perfume.jpg",          8.5, 2.2, 2.0, 2.2, -6)
img("tattoo.jpg",           10.5, 3.2, 2.5, 1.9, 4)
img("mirror.jpg",           12.8, 2.8, 2.2, 1.8, -3)
img("chocolate.jpg",        13.0, 4.4, 2.5, 1.9, 7)

# --- Cluster onder: WILDE WERELD ---
img("octopus.jpg",          2.5, 5.8, 2.3, 1.8, -5)
img("parrot.jpg",           4.5, 6.0, 2.0, 1.6, 8)
img("jellyfish.jpg",        6.3, 5.6, 2.2, 1.8, -3)
img("butterfly.jpg",        10.2, 5.5, 2.5, 1.9, 4)
img("yoga.jpg",             12.3, 5.8, 2.3, 1.7, -6)

# --- Onderste strip: DROMEN ---
img("neon-city.jpg",        0.0, 7.2, 2.8, 1.7, 3)
img("graffiti.jpg",         2.5, 7.0, 2.5, 1.8, -4)
img("fire.jpg",             4.7, 7.3, 2.2, 1.6, 5)
img("sand-writing.jpg",     6.6, 7.0, 2.5, 1.7, -2)
img("tree-rings.jpg",       8.8, 7.3, 2.3, 1.5, 6)
img("maze.jpg",             10.8, 7.0, 2.5, 1.8, -5)
img("space.jpg",            12.9, 7.2, 3.0, 1.7, 3)
img("cooking.jpg",          14.5, 5.5, 1.5, 1.3, -8)
img("wine-cork.jpg",        14.3, 3.5, 1.6, 1.4, 10)

# Extra overlaps
img("gold-texture.jpg",     3.5, 1.5, 1.8, 1.3, 12)
img("leather.jpg",          7.2, 5.8, 1.5, 1.2, -10)
img("metal-texture.jpg",    5.8, 7.8, 1.6, 1.0, 8)
img("ink-paper.jpg",        8.5, 6.8, 1.8, 1.4, -7)
img("lego.jpg",             14.0, 0.3, 1.8, 1.5, -12)

# ═══════════════════════════════════════════════════
# CHAOTISCHE TAGS — gedraaid, verspreid, alle kanten op
# ═══════════════════════════════════════════════════

# --- Vormen (niet alleen een pen!) ---
tag(6.6, 0.0, "RING MET SENSOR?", CYAN, 9, 1.8, -3)
tag(7.0, 2.0, "SMARTWATCH = PEN-VERVANGER?", PINK, 8, 2.6, 5)
tag(10.3, 1.5, "HORLOGE DAT JE SCHRIFT HERKENT", LIME, 7, 2.8, -4)
tag(13.5, 2.2, "SPIEGEL DIE JE SCANT", PURPLE, 8, 2.0, 8)
tag(11.0, 4.8, "TATTOO ALS PASPOORT", ORANGE, 8, 2.0, -6)
tag(14.2, 3.2, "WIJN PROEVEN = 2FA", RED, 7, 1.6, 12)
tag(13.2, 6.0, "YOGA-POSE = UNLOCK", PURPLE, 7, 1.8, -8)

# --- Authenticatie methodes ---
tag(0.1, 0.0, "IRIS", CYAN, 12, 0.8, -5)
tag(3.0, 2.3, "VINGERAFDRUK", PINK, 9, 1.6, 7)
tag(4.8, 2.0, "SCHRIJF-DNA", GOLD, 10, 1.4, -3)
tag(0.2, 3.5, "LAKZEGEL 2.0", ORANGE, 9, 1.4, 4)
tag(2.5, 3.2, "KLOP-RITME = CODE", LIME, 8, 2.0, -6)
tag(8.0, 4.5, "CHIP IN DE PUNT", CYAN, 8, 1.6, 3)

# --- Productvorm-ideeën ---
tag(0.3, 5.2, "KASTJE NAAST JE PC", YELLOW, 8, 2.0, -4)
tag(5.0, 5.0, "IN JE TELEFOONHOESJE", PINK, 8, 2.0, 6)
tag(6.3, 3.5, "INGEBOUWD IN LAPTOP", LIME, 8, 2.0, -3)
tag(8.5, 3.0, "PARFUM DAT JE AUTH", PURPLE, 8, 2.0, 5)
tag(3.0, 6.5, "HANDSCHOEN MET SENSOREN", CYAN, 7, 2.3, -5)
tag(6.8, 6.5, "BUREAUMAT DIE MEELEEST", ORANGE, 7, 2.3, 4)
tag(9.5, 6.0, "STEMPEL (maar slim)", GOLD, 8, 1.8, -7)

# --- Grappige referenties ---
tag(0.5, 7.0, "BANKSY = ANONIEM\nMAAR HERKENBAAR", PINK, 6, 2.0, 5)
tag(5.0, 7.0, "VUUR-LOPEN ALS LOGIN", RED, 7, 1.8, -6)
tag(9.0, 7.0, "JAARRINGEN = JE LEEFTIJD\nIN JE SCHRIJFSTIJL", LIME, 6, 2.2, 3)
tag(11.0, 7.0, "DOOLHOF: ALLEEN JIJ\nKENT DE WEG", PURPLE, 7, 2.2, -4)
tag(13.2, 7.0, "ONDERTEKENEN\nVANUIT DE RUIMTE", CYAN, 7, 2.0, 7)

# ═══════════════════════════════════════════════════
# WILDE VRAGEN — schots en scheef
# ═══════════════════════════════════════════════════

vragen = [
    (0.4, 1.8, "Wat als de PEN\nweg kan?", 9, GOLD, -8),
    (4.8, 0.2, "Wat als het in\nje HORLOGE zit?", 8, CYAN, 6),
    (11.5, 0.0, "Wat als je DANS\nje wachtwoord is?", 8, PINK, -4),
    (2.0, 4.5, "Wat als de DEUR\nje handschrift kent?", 7, LIME, 5),
    (5.5, 5.5, "Wat als papier\nzelf weet WIE\nerop schrijft?", 7, ORANGE, -6),
    (10.0, 5.0, "Vlinder = uniek.\nSneeuwvlok = uniek.\nSchrijfstijl = ???", 7, CYAN, 3),
    (0.3, 6.5, "Octopus: 8 armen\n8 sensoren\n8 wachtwoorden", 7, PURPLE, -3),
    (4.6, 6.3, "Als papegaai je\nstem nadoet:\nis hij jou?", 7, PINK, 7),
    (7.5, 7.5, "Oma's geheime recept =\noude 2FA", 7, GOLD, -5),
    (14.0, 1.0, "LEGO: bouw je\neigen authenticatie\nper situatie", 8, LIME, -15),
    (14.0, 6.5, "Koken op gevoel.\nSchrijven op gevoel.\nZelfde ding?", 6, ORANGE, 10),
    (9.8, 2.0, "GEUR is de meest\npersoonlijke zin.\nParfum = handtekening?", 7, PURPLE, -3),
]
for x, y, text, sz, col, rot in vragen:
    lbl(x, y, text, sz, col, 2.0, rot)

# ═══════════════════════════════════════════════════
# ACCENT DOTS — energie, chaos
# ═══════════════════════════════════════════════════
dots = [
    (1.5, 0.3, 0.15, PINK), (5.2, 1.2, 0.12, LIME), (7.8, 0.5, 0.18, GOLD),
    (10.0, 0.8, 0.1, CYAN), (13.8, 1.8, 0.14, ORANGE), (3.8, 3.0, 0.16, PURPLE),
    (9.5, 3.5, 0.12, PINK), (11.8, 5.2, 0.15, LIME), (1.0, 7.5, 0.13, GOLD),
    (7.0, 8.0, 0.1, CYAN), (15.0, 7.0, 0.16, PINK), (12.0, 0.5, 0.11, RED),
    (4.0, 7.8, 0.14, YELLOW), (14.5, 5.0, 0.12, PURPLE), (0.8, 4.5, 0.1, LIME),
]
for x, y, sz, col in dots:
    circle(x, y, sz, col)

# ═══════════════════════════════════════════════════
# CENTRAAL LABEL (klein, niet dominant)
# ═══════════════════════════════════════════════════
c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.8), Inches(4.0), Inches(0.35))
c.fill.solid(); c.fill.fore_color.rgb = DARK; c.line.color.rgb = GOLD; c.line.width = Pt(2)
tf = c.text_frame; p = tf.paragraphs[0]
p.text = "HOE BEWIJS JE WIE JE BENT?"; p.font.size = Pt(13)
p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = "Georgia"; p.alignment = PP_ALIGN.CENTER

# Subtitel
c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(6.2), Inches(5.0), Inches(0.28))
c2.fill.solid(); c2.fill.fore_color.rgb = DARK; c2.line.fill.background()
tf = c2.text_frame; p = tf.paragraphs[0]
p.text = "Niet per se een pen. Misschien een ring. Of je huid. Of je dans."; p.font.size = Pt(9)
p.font.color.rgb = CYAN; p.font.bold = True; p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER

# Title (top, small)
t = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(8.7), Inches(6.0), Inches(0.25))
t.fill.solid(); t.fill.fore_color.rgb = GOLD; t.line.fill.background()
tf = t.text_frame; p = tf.paragraphs[0]
p.text = "CONCEPT COLLECTION — DE ANALOGE SECURITY EXPERT — DIVERGEER!"; p.font.size = Pt(8)
p.font.color.rgb = DARK; p.font.bold = True; p.font.name = "Georgia"; p.alignment = PP_ALIGN.CENTER

output = "/Users/jr/Library/CloudStorage/OneDrive-Zonneplan/School/Innovatiemanagement/public/concept-collection-security-expert.pptx"
prs.save(output)
print(f"Saved: {output}")
